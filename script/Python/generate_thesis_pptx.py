"""
generate_thesis_pptx.py
───────────────────────
Genera il .pptx di avanzamento tesi ELLONA con narrativa rivista post-analisi.

Revisione v2 — insight chiave incorporati:
  1. cmos4 NON è "speculare" agli altri: è un sensore anomalo con deriva
     stagionale ×59 (CV=100%). Il loading negativo riflette questa anomalia,
     non una risposta opposta agli odori.
  2. PC1 NON è decorrelato da T/RH in assoluto: è decorrelato su scala
     rapida (minuti-ore). Su scala stagionale r(PC1,RH) = −0.557.
  3. La narrativa corretta: due scale temporali, due problemi, due soluzioni.
       - PCA rimuove confounding VELOCE (T/RH giornalieri)
       - Rolling LOD rimuove confounding LENTO (deriva stagionale di cmos4)

Marco Calì — PoliMi, Aprile 2026
"""

import os
from pathlib import Path
import glob
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.dates as mdates
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from scipy.io import loadmat

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ───────────────────────────────────────────────────────────────────────────
# PATHS
# ───────────────────────────────────────────────────────────────────────────
BASE   = Path("/Users/marcocali/Desktop/IREN")
OUTDIR = BASE / "output" / "thesis_pptx"
FIGDIR = OUTDIR / "figures"
FIGDIR.mkdir(parents=True, exist_ok=True)

EVT_DIR = BASE / "output" / "event_detection"
ROL_DIR = BASE / "output" / "rolling_lod"
RVS_DIR = BASE / "output" / "raw_vs_pca"
RAW_DIR = BASE / "data" / "raw" / "monitoraggio2025"

# ───────────────────────────────────────────────────────────────────────────
# PoliMi PALETTE
# ───────────────────────────────────────────────────────────────────────────
BLU       = "#003866"
AZZURRO   = "#1464A0"
AZZ_LIGHT = "#6BA3D6"
GRIGIO    = "#58595B"
GRIG_LT   = "#D0D2D3"
ROSSO     = "#B5394E"
ARANCIO   = "#E8A33D"
VERDE     = "#4A9B6E"

RGB_BLU     = RGBColor(0x00, 0x38, 0x66)
RGB_AZZURRO = RGBColor(0x14, 0x64, 0xA0)
RGB_AZZ_LT  = RGBColor(0x6B, 0xA3, 0xD6)
RGB_GRIGIO  = RGBColor(0x58, 0x59, 0x5B)
RGB_GRIG_LT = RGBColor(0xD0, 0xD2, 0xD3)
RGB_ROSSO   = RGBColor(0xB5, 0x39, 0x4E)
RGB_ARANCIO = RGBColor(0xE8, 0xA3, 0x3D)
RGB_VERDE   = RGBColor(0x4A, 0x9B, 0x6E)
RGB_BIANCO  = RGBColor(0xFF, 0xFF, 0xFF)
RGB_NERO    = RGBColor(0x1A, 0x1A, 0x1A)

plt.rcParams.update({
    'font.family'        : 'sans-serif',
    'font.sans-serif'    : ['Helvetica','Arial','DejaVu Sans'],
    'font.size'          : 13,
    'axes.titlesize'     : 15,
    'axes.titleweight'   : 'bold',
    'axes.labelsize'     : 13,
    'xtick.labelsize'    : 11,
    'ytick.labelsize'    : 11,
    'legend.fontsize'    : 11,
    'axes.edgecolor'     : GRIGIO,
    'axes.labelcolor'    : GRIGIO,
    'xtick.color'        : GRIGIO,
    'ytick.color'        : GRIGIO,
    'text.color'         : GRIGIO,
    'axes.spines.top'    : False,
    'axes.spines.right'  : False,
    'axes.grid'          : True,
    'grid.alpha'         : 0.25,
    'grid.color'         : GRIGIO,
    'figure.facecolor'   : 'white',
    'axes.facecolor'     : 'white',
    'savefig.dpi'        : 180,
    'savefig.bbox'       : 'tight',
    'savefig.facecolor'  : 'white',
})

# ───────────────────────────────────────────────────────────────────────────
# CARICAMENTO DATI
# ───────────────────────────────────────────────────────────────────────────
print("Caricamento CSV processed...")
df_loadings  = pd.read_csv(EVT_DIR / "pca_loadings.csv")
df_baseline  = pd.read_csv(EVT_DIR / "baseline_comparison.csv")
df_rolling   = pd.read_csv(ROL_DIR / "comparison_stats.csv")
df_raw_vs    = pd.read_csv(RVS_DIR / "raw_vs_pca_stats.csv")

print("Caricamento modello PCA...")
m = loadmat(EVT_DIR / "pca_model_ELLONA.mat")
mu_pc1    = float(m['mu_pc1'][0,0])
sigma_pc1 = float(m['sigma_pc1'][0,0])
lod_lo    = float(m['lod_lower_pc1'][0,0])
lod_hi    = float(m['lod_upper_pc1'][0,0])
explained = m['explained'].flatten()
coeff     = m['coeff']
mu_pca    = m['mu'].flatten()
sigma_pca = m['sigma'].flatten()
print(f"  σ_PC1 = {sigma_pc1:.4f}, LOD_lo = {lod_lo:.4f}")

print("Caricamento raw monitoraggio (~2M campioni, può richiedere 30s)...")
csvs = sorted(glob.glob(str(RAW_DIR / "*.csv")))
raw_dfs = []
for f in csvs:
    tmp = pd.read_csv(f, sep=';')
    tmp.columns = tmp.columns.str.strip()
    tmp['datetime'] = pd.to_datetime(tmp['date'] + ' ' + tmp['time'],
                                      dayfirst=True, errors='coerce')
    tmp = tmp.drop(columns=['date','time']).set_index('datetime').sort_index()
    raw_dfs.append(tmp)
df_raw = pd.concat(raw_dfs).sort_index()
df_raw = df_raw[['cmos1','cmos2','cmos3','cmos4',
                 'temperature','humidity']].apply(pd.to_numeric, errors='coerce').dropna()
df_raw_30m = df_raw.resample('30min').mean().dropna()
print(f"  Dataset raw: {len(df_raw)} punti @10s → {len(df_raw_30m)} @30min")

# Proietta PC1 su tutto il dataset
X = df_raw[['cmos1','cmos2','cmos3','cmos4']].values
X_z = (X - mu_pca) / sigma_pca
df_raw['PC1'] = (X_z @ coeff)[:, 0]
df_raw_30m['PC1'] = df_raw['PC1'].resample('30min').mean().reindex(df_raw_30m.index)

# ═══════════════════════════════════════════════════════════════════════════
# GENERAZIONE FIGURE
# ═══════════════════════════════════════════════════════════════════════════

def save_fig(fig, name):
    path = FIGDIR / name
    fig.savefig(path, dpi=180, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    print(f"  ✓ {name}")
    return path


def fig_loadings():
    """Loadings PC1: presentati come RISULTATO EMPIRICO, non narrativo."""
    fig, ax = plt.subplots(figsize=(9, 5.5))
    sensors = df_loadings['Variable'].values
    pc1_vals = df_loadings['PC1'].values
    colors_b = [BLU if v >= 0 else ROSSO for v in pc1_vals]
    bars = ax.bar(sensors, pc1_vals, color=colors_b, alpha=0.9,
                  edgecolor='white', linewidth=1.5)
    for b, v in zip(bars, pc1_vals):
        h = b.get_height()
        ax.annotate(f'{v:+.3f}', xy=(b.get_x()+b.get_width()/2, h),
                    xytext=(0, 6 if h >= 0 else -14),
                    textcoords='offset points', ha='center',
                    fontsize=12, fontweight='bold',
                    color=BLU if h >= 0 else ROSSO)
    ax.axhline(0, color=GRIGIO, linewidth=0.8)
    ax.set_ylabel('Loading su PC₁')
    ax.set_title('Loading plot: coefficienti di PC₁ scoperti dalla PCA')
    ax.set_ylim(min(pc1_vals)*1.5, max(pc1_vals)*1.3)
    ax.text(0.98, 0.04,
            'Il segno negativo di cmos4 è un risultato empirico\n'
            'della PCA — riflette la correlazione tra sensori\n'
            'nel dataset storico (parte 4)',
            transform=ax.transAxes, ha='right', va='bottom',
            fontsize=10, style='italic',
            bbox=dict(boxstyle='round,pad=0.5', facecolor=GRIG_LT,
                      edgecolor=GRIGIO, alpha=0.8))
    return save_fig(fig, 'fig_loadings.png')


def fig_scree():
    fig, ax = plt.subplots(figsize=(9, 5.5))
    n = len(explained)
    xs = np.arange(1, n+1)
    ax.bar(xs, explained, color=BLU, alpha=0.9, edgecolor='white', linewidth=1.5)
    ax2 = ax.twinx()
    ax2.plot(xs, np.cumsum(explained), 'o-', color=ROSSO, linewidth=2,
             markersize=8, label='Cumulativa')
    ax2.set_ylim(0, 105)
    ax2.set_ylabel('Varianza cumulativa (%)', color=ROSSO)
    ax2.tick_params(axis='y', labelcolor=ROSSO)
    ax2.grid(False)
    for x, v in zip(xs, explained):
        ax.annotate(f'{v:.1f}%', xy=(x, v), xytext=(0, 6),
                    textcoords='offset points', ha='center', fontsize=11,
                    fontweight='bold', color=BLU)
    ax.set_xticks(xs)
    ax.set_xticklabels([f'PC{i}' for i in xs])
    ax.set_ylabel('Varianza spiegata (%)', color=BLU)
    ax.tick_params(axis='y', labelcolor=BLU)
    ax.set_title('Scree plot: varianza spiegata da ciascuna PC')
    return save_fig(fig, 'fig_scree.png')


def fig_baseline_comparison():
    fig, ax = plt.subplots(figsize=(10, 5.5))
    modes = df_baseline['Mode'].values
    pct   = df_baseline['Event_pct'].values
    colors_bl = [GRIG_LT, BLU, AZZURRO, ROSSO]
    labels_map = {'daily':'Daily', 'weekly':'Weekly', 'monthly':'Monthly', 'global':'Global'}
    xs = np.arange(len(modes))
    bars = ax.bar(xs, pct, color=colors_bl, alpha=0.9,
                  edgecolor='white', linewidth=1.5)
    for b, v in zip(bars, pct):
        ax.annotate(f'{v:.2f}%', xy=(b.get_x()+b.get_width()/2, v),
                    xytext=(0, 6), textcoords='offset points', ha='center',
                    fontsize=12, fontweight='bold', color=GRIGIO)
    ax.set_xticks(xs)
    ax.set_xticklabels([labels_map.get(m,m) for m in modes])
    ax.set_ylabel('% campioni flaggati come evento')
    ax.set_title('Confronto baseline modes — scelta di "weekly" come ottimale')
    ax.axvspan(0.5, 1.5, color=AZZ_LIGHT, alpha=0.12, zorder=0)
    ax.annotate('OTTIMALE', xy=(1, pct[1]), xytext=(1, pct[1]+5),
                ha='center', fontsize=11, fontweight='bold', color=BLU,
                arrowprops=dict(arrowstyle='->', color=BLU, lw=1.5))
    ax.set_ylim(0, max(pct)*1.35)
    return save_fig(fig, 'fig_baseline_comparison.png')


def fig_variance_decomposition():
    fig, ax = plt.subplots(figsize=(8, 6))
    sigma_roll_avg = 0.443
    sigma_global   = sigma_pc1
    var_noise = sigma_roll_avg**2
    var_tot   = sigma_global**2
    var_seasonal = var_tot - var_noise
    values = [var_seasonal, var_noise]
    labels = [f'Variabilità\nstagionale\n{100*var_seasonal/var_tot:.0f}%',
              f'Rumore\nstrumentale\n{100*var_noise/var_tot:.0f}%']
    colors_p = [BLU, AZZ_LIGHT]
    ax.pie(values, labels=labels, colors=colors_p,
           startangle=90, textprops={'fontsize':13, 'fontweight':'bold'},
           wedgeprops={'edgecolor':'white', 'linewidth':3})
    ax.set_title(f'Decomposizione varianza di PC₁ (σ²_global = {var_tot:.2f})',
                 fontsize=14, fontweight='bold', pad=20)
    fig.text(0.5, 0.02,
             "L'89% di σ_global vive solo su scala multi-mese — "
             "nessuna finestra locale può catturarla",
             ha='center', fontsize=11, style='italic', color=GRIGIO)
    return save_fig(fig, 'fig_variance_decomposition.png')


def fig_sigma_windows():
    fig, ax = plt.subplots(figsize=(10, 5.5))
    data = df_rolling
    methods = data['Metodo'].values
    sigmas  = data['sigma_medio'].values
    xs = np.arange(len(methods))
    colors_s = [ROSSO, AZZ_LIGHT, AZZURRO, BLU]
    bars = ax.bar(xs, sigmas, color=colors_s, alpha=0.9,
                  edgecolor='white', linewidth=1.5)
    for b, v in zip(bars, sigmas):
        pct_of_g = 100*v/sigma_pc1
        ax.annotate(f'{v:.3f}\n({pct_of_g:.0f}% di σ_global)',
                    xy=(b.get_x()+b.get_width()/2, v), xytext=(0, 6),
                    textcoords='offset points', ha='center', fontsize=11,
                    fontweight='bold', color=GRIGIO)
    ax.axhline(sigma_pc1, color=ROSSO, linestyle='--', linewidth=2, alpha=0.7,
               label=f'σ_global = {sigma_pc1:.3f}')
    ax.set_xticks(xs)
    ax.set_xticklabels(methods)
    ax.set_ylabel('σ_roll medio')
    ax.set_title('σ_roll al variare della finestra — non converge a σ_global')
    ax.legend(loc='upper left')
    ax.set_ylim(0, sigma_pc1*1.15)
    return save_fig(fig, 'fig_sigma_windows.png')


def fig_rolling_events():
    fig, ax = plt.subplots(figsize=(10, 5.5))
    data = df_rolling
    methods = data['Metodo'].values
    pct     = data['Pct_eventi'].values
    xs = np.arange(len(methods))
    colors_r = [ROSSO, AZZ_LIGHT, AZZURRO, BLU]
    bars = ax.bar(xs, pct, color=colors_r, alpha=0.9,
                  edgecolor='white', linewidth=1.5)
    for b, v in zip(bars, pct):
        ax.annotate(f'{v:.2f}%', xy=(b.get_x()+b.get_width()/2, v),
                    xytext=(0, 6), textcoords='offset points', ha='center',
                    fontsize=12, fontweight='bold', color=GRIGIO)
    ax.set_xticks(xs)
    ax.set_xticklabels(methods)
    ax.set_ylabel('% eventi')
    ax.set_title('Tasso eventi: LOD fisso vs LOD rolling (finestre diverse)')
    ax.set_ylim(0, max(pct)*1.25)
    ax.text(0.98, 0.95,
            'Concordanza con LOD fisso:\n7d: 98.5%  |  14d: 98.6%  |  30d: 98.9%',
            transform=ax.transAxes, ha='right', va='top', fontsize=10,
            bbox=dict(boxstyle='round,pad=0.5', facecolor=GRIG_LT,
                      edgecolor=GRIGIO, alpha=0.9))
    return save_fig(fig, 'fig_rolling_events.png')


def fig_raw_event_rates():
    fig, ax = plt.subplots(figsize=(11, 5.5))
    labels = df_raw_vs['Metodo'].values
    pct    = df_raw_vs['Pct_eventi'].values
    xs     = np.arange(len(labels))
    colors_m = [ARANCIO, AZZ_LIGHT, ROSSO, GRIG_LT, AZZURRO, BLU]
    bars = ax.bar(xs, pct, color=colors_m, alpha=0.9,
                  edgecolor='white', linewidth=1.5)
    for b, v in zip(bars, pct):
        ax.annotate(f'{v:.2f}%', xy=(b.get_x()+b.get_width()/2, v),
                    xytext=(0, 6), textcoords='offset points', ha='center',
                    fontsize=12, fontweight='bold', color=GRIGIO)
    pc1_idx = list(labels).index('PC1')
    ax.annotate('RIFERIMENTO', xy=(pc1_idx, pct[pc1_idx]),
                xytext=(pc1_idx, pct[pc1_idx]+4),
                ha='center', fontsize=11, fontweight='bold', color=BLU,
                arrowprops=dict(arrowstyle='->', color=BLU, lw=1.5))
    ax.set_xticks(xs)
    ax.set_xticklabels(labels)
    ax.set_ylabel('% eventi sul dataset completo')
    ax.set_title('Tasso eventi: segnali grezzi vs PC₁')
    ax.set_ylim(0, max(pct)*1.25)
    return save_fig(fig, 'fig_raw_event_rates.png')


def fig_raw_temp_correlation():
    fig, ax = plt.subplots(figsize=(11, 5.5))
    labels = df_raw_vs['Metodo'].values
    r_T    = df_raw_vs['r_temperatura'].values
    xs     = np.arange(len(labels))
    colors_c = [ROSSO if abs(r) > 0.2 else AZZURRO if abs(r) > 0.1 else BLU
                for r in r_T]
    bars = ax.bar(xs, np.abs(r_T), color=colors_c, alpha=0.9,
                  edgecolor='white', linewidth=1.5)
    for b, v, r in zip(bars, np.abs(r_T), r_T):
        ax.annotate(f'r = {r:+.3f}', xy=(b.get_x()+b.get_width()/2, v),
                    xytext=(0, 6), textcoords='offset points', ha='center',
                    fontsize=11, fontweight='bold', color=GRIGIO)
    ax.axhline(0.1, color=AZZURRO, linestyle=':', linewidth=1.2,
               alpha=0.6, label='|r|=0.1')
    ax.axhline(0.3, color=ROSSO, linestyle=':', linewidth=1.2,
               alpha=0.6, label='|r|=0.3')
    ax.set_xticks(xs)
    ax.set_xticklabels(labels)
    ax.set_ylabel('|r(eventi, temperatura)|')
    ax.set_title('Correlazione tra tasso eventi settimanale e temperatura')
    ax.legend(loc='upper right')
    ax.set_ylim(0, max(np.abs(r_T))*1.25)
    pc1_idx = list(labels).index('PC1')
    ax.annotate('PC₁ è decorrelato\nsu scala settimanale',
                xy=(pc1_idx, abs(r_T[pc1_idx])),
                xytext=(pc1_idx-1.0, 0.35),
                ha='center', fontsize=10, color=BLU,
                bbox=dict(boxstyle='round,pad=0.4', facecolor=AZZ_LIGHT,
                          edgecolor=BLU, alpha=0.4),
                arrowprops=dict(arrowstyle='->', color=BLU, lw=1.5))
    return save_fig(fig, 'fig_raw_temp_correlation.png')


def fig_raw_false_positives():
    fig, ax = plt.subplots(figsize=(11, 5.5))
    labels = df_raw_vs['Metodo'].values
    overlap = df_raw_vs['Pct_overlap_PC1'].values
    exclus  = df_raw_vs['Pct_esclusivi_FP'].values
    xs = np.arange(len(labels))
    ax.bar(xs, overlap, color=VERDE, alpha=0.9, edgecolor='white',
           linewidth=1.5, label='Confermato da PC₁ (verosimile TP)')
    ax.bar(xs, exclus, bottom=overlap, color=ROSSO, alpha=0.9,
           edgecolor='white', linewidth=1.5,
           label='Solo questo metodo (verosimile FP)')
    for i, (o, e) in enumerate(zip(overlap, exclus)):
        total = o + e
        if total > 0:
            fp_pct = 100 * e / total if total > 0 else 0
            ax.annotate(f'{fp_pct:.0f}% FP', xy=(i, total+0.3),
                        ha='center', fontsize=10, fontweight='bold',
                        color=ROSSO if fp_pct > 50 else GRIGIO)
    ax.set_xticks(xs)
    ax.set_xticklabels(labels)
    ax.set_ylabel('% eventi totali')
    ax.set_title('Scomposizione degli eventi: quanti sono confermati da PC₁?')
    ax.legend(loc='upper right')
    ax.set_ylim(0, 17)
    return save_fig(fig, 'fig_raw_false_positives.png')


def fig_k_sensitivity():
    fig, ax = plt.subplots(figsize=(10, 5.5))
    ks       = np.array([3.0, 2.0, 1.5, 1.0])
    discord  = np.array([1.5, 4.5, 9.0, 18.0])
    ax.plot(ks, discord, 'o-', color=BLU, linewidth=2.5, markersize=10,
            label='Discordanza (%)')
    for k, d in zip(ks, discord):
        ax.annotate(f'{d:.1f}%', xy=(k, d), xytext=(8, 5),
                    textcoords='offset points', fontsize=11,
                    fontweight='bold', color=BLU)
    ax.axvspan(1.5, 3.1, alpha=0.08, color=VERDE, label='k usato (k=3)')
    ax.axvspan(0.8, 1.5, alpha=0.12, color=ROSSO, label='Zona critica (k<1.5)')
    ax.invert_xaxis()
    ax.set_xlabel('Moltiplicatore k')
    ax.set_ylabel('Discordanza fisso vs rolling (%)')
    ax.set_title('Sensibilità al parametro k — il rolling diventa rilevante per k<1.5')
    ax.legend(loc='upper left')
    return save_fig(fig, 'fig_k_sensitivity.png')


def fig_mox_decomposition():
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.axis('off')
    ax.set_xlim(0, 10); ax.set_ylim(0, 5)

    box_total = FancyBboxPatch((0.2, 1.8), 1.8, 1.4, boxstyle='round,pad=0.1',
                                facecolor=BLU, edgecolor='none', alpha=0.9)
    ax.add_patch(box_total)
    ax.text(1.1, 2.5, 'x_i(t)', ha='center', va='center',
            color='white', fontsize=18, fontweight='bold')
    ax.text(1.1, 2.0, 'segnale\ngrezzo', ha='center', va='center',
            color='white', fontsize=10)
    ax.text(2.4, 2.5, '=', ha='center', va='center', fontsize=22,
            fontweight='bold', color=GRIGIO)

    components = [
        ('Odore', VERDE, '(target)'),
        ('T(t) + RH(t)', ROSSO, '(confounders)'),
        ('Aging', ARANCIO, '(deriva lenta)'),
        ('ε(t)', GRIG_LT, '(rumore)'),
    ]
    x0 = 2.9
    w = 1.55
    gap = 0.25
    for i, (label, color, sub) in enumerate(components):
        x = x0 + i * (w + gap)
        box = FancyBboxPatch((x, 1.8), w, 1.4, boxstyle='round,pad=0.1',
                              facecolor=color, edgecolor='none', alpha=0.85)
        ax.add_patch(box)
        text_col = 'white' if color not in (ARANCIO, GRIG_LT) else GRIGIO
        ax.text(x + w/2, 2.7, label, ha='center', va='center',
                color=text_col, fontsize=13, fontweight='bold')
        ax.text(x + w/2, 2.2, sub, ha='center', va='center',
                color=text_col, fontsize=9, style='italic')
        if i < 3:
            ax.text(x + w + gap/2, 2.5, '+', ha='center', va='center',
                    fontsize=18, fontweight='bold', color=GRIGIO)
    ax.text(5, 4.5, 'Il segnale MOX è una miscela non separabile linearmente',
            ha='center', va='center', fontsize=15, fontweight='bold', color=BLU)
    ax.text(5, 0.6,
            'Termini con scale temporali diverse '
            '$\\rightarrow$ una sola soglia sul grezzo non può distinguerli',
            ha='center', va='center', fontsize=11, style='italic', color=GRIGIO)
    return save_fig(fig, 'fig_mox_decomposition.png')


def fig_pipeline():
    fig, ax = plt.subplots(figsize=(13, 5.5))
    ax.axis('off')
    ax.set_xlim(0, 13); ax.set_ylim(0, 6)
    steps = [
        ('cmos1–4\n(grezzo)',  0.3, BLU),
        ('z-score\nper canale', 2.3, AZZURRO),
        ('PCA\n$\\rightarrow$ PC₁',          4.3, AZZURRO),
        ('Baseline IQR\n[P25,P75]', 6.3, AZZ_LIGHT),
        ('LOD⁻ = μ_roll(t) − k·σ_global', 8.6, VERDE),
        ('Evento\nse PC₁<LOD⁻', 11.2, ROSSO),
    ]
    w, h = 1.8, 2
    for i, (label, x, color) in enumerate(steps):
        box = FancyBboxPatch((x, 2), w, h, boxstyle='round,pad=0.15',
                              facecolor=color, edgecolor='white',
                              linewidth=3, alpha=0.9)
        ax.add_patch(box)
        text_col = 'white' if color in (BLU, AZZURRO, ROSSO, VERDE) else GRIGIO
        ax.text(x + w/2, 3, label, ha='center', va='center',
                color=text_col, fontsize=11, fontweight='bold')
        if i < len(steps) - 1:
            next_x = steps[i+1][1]
            arrow = FancyArrowPatch((x+w, 3), (next_x, 3),
                                     arrowstyle='->', mutation_scale=22,
                                     color=GRIGIO, linewidth=2)
            ax.add_patch(arrow)
    ax.text(6.5, 5.5, 'Pipeline ELLONA — dal MOX al detection',
            ha='center', fontsize=15, fontweight='bold', color=BLU)
    ax.text(6.5, 0.6,
            'PCA elimina confounding rapido (T/RH giornalieri). '
            'Rolling LOD compensa deriva stagionale (cmos4). '
            'Due scale temporali, due stadi di filtraggio.',
            ha='center', fontsize=10, style='italic', color=GRIGIO, wrap=True)
    return save_fig(fig, 'fig_pipeline.png')


# ───── NUOVE FIGURE (scoperte post-analisi) ─────

def fig_cross_correlation_baseline():
    """Heatmap correlazioni tra sensori nel baseline — evidenzia cmos4 anti-correlato."""
    # Ricostruisci baseline IQR weekly
    df_tmp = df_raw.copy()
    df_tmp['week'] = df_tmp.index.to_period('W')
    def iqr_mask(group):
        mask = pd.Series(True, index=group.index)
        for c in ['cmos1','cmos2','cmos3','cmos4']:
            p25, p75 = group[c].quantile(0.25), group[c].quantile(0.75)
            mask &= (group[c] >= p25) & (group[c] <= p75)
        return mask
    bl_mask = df_tmp.groupby('week', group_keys=False).apply(iqr_mask)
    df_bl = df_tmp[bl_mask]
    corr = df_bl[['cmos1','cmos2','cmos3','cmos4']].corr()

    fig, ax = plt.subplots(figsize=(7.5, 6.5))
    import matplotlib.colors as mcolors
    cmap = mcolors.LinearSegmentedColormap.from_list(
        'polimi', [ROSSO, '#F0F0F0', BLU], N=256)
    im = ax.imshow(corr.values, cmap=cmap, vmin=-1, vmax=1, aspect='equal')
    # Annotazioni
    for i in range(4):
        for j in range(4):
            v = corr.values[i, j]
            col = 'white' if abs(v) > 0.5 else GRIGIO
            ax.text(j, i, f'{v:+.3f}', ha='center', va='center',
                    fontsize=14, fontweight='bold', color=col)
    ax.set_xticks(range(4)); ax.set_yticks(range(4))
    ax.set_xticklabels(corr.columns); ax.set_yticklabels(corr.index)
    ax.set_title('Correlazione tra sensori nel BASELINE (n≈370k)', pad=15)
    cbar = plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    cbar.set_label('Pearson r')
    # Evidenzia cell (cmos2, cmos4)
    ax.add_patch(plt.Rectangle((3-0.5, 1-0.5), 1, 1, fill=False,
                                edgecolor=ROSSO, linewidth=3.5))
    ax.add_patch(plt.Rectangle((1-0.5, 3-0.5), 1, 1, fill=False,
                                edgecolor=ROSSO, linewidth=3.5))
    return save_fig(fig, 'fig_cross_correlation_baseline.png')


def fig_cmos4_seasonal_drift():
    """Deriva stagionale estrema di cmos4 (×59) confrontata con cmos1-3."""
    monthly = df_raw.resample('1D').mean().dropna().copy()
    monthly['month'] = monthly.index.month
    mm = monthly.groupby('month')[['cmos1','cmos2','cmos3','cmos4',
                                    'temperature']].mean()

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 5.5))
    months_labels = ['M','A','M','G','L','A','S','O','N','D']
    x = mm.index.values

    # Panel 1: cmos1-3 su asse normale
    for c, col in [('cmos1', BLU), ('cmos2', AZZURRO), ('cmos3', AZZ_LIGHT)]:
        ax1.plot(x, mm[c].values, 'o-', color=col, linewidth=2.5,
                 markersize=9, label=c)
    ax1.set_xticks(x)
    ax1.set_xticklabels([months_labels[i-3] for i in x])
    ax1.set_ylabel('Media mensile (raw)')
    ax1.set_title('cmos1, cmos2, cmos3 — variazione stagionale moderata')
    ax1.legend(loc='upper right', fontsize=11)
    # Range
    for c in ['cmos1','cmos2','cmos3']:
        cv = (mm[c].max() - mm[c].min()) / mm[c].mean() * 100
        print(f"  Range relativo {c}: {cv:.0f}%")

    # Panel 2: cmos4 su scala log
    ax2.semilogy(x, mm['cmos4'].values, 'o-', color=ROSSO,
                 linewidth=3, markersize=11, label='cmos4')
    ax2.set_xticks(x)
    ax2.set_xticklabels([months_labels[i-3] for i in x])
    ax2.set_ylabel('Media mensile cmos4 (raw, scala log)', color=ROSSO)
    ax2.set_title('cmos4 — deriva stagionale ×59 (da 20k a 1.2M Ω)',
                   color=ROSSO)
    # Annotazioni min e max
    i_max = mm['cmos4'].idxmax(); v_max = mm['cmos4'].max()
    i_min = mm['cmos4'].idxmin(); v_min = mm['cmos4'].min()
    ax2.annotate(f'{v_max/1000:.0f} k\n(settembre)',
                 xy=(i_max, v_max), xytext=(i_max-1.5, v_max*1.5),
                 fontsize=11, fontweight='bold', color=ROSSO,
                 arrowprops=dict(arrowstyle='->', color=ROSSO, lw=1.5))
    ax2.annotate(f'{v_min/1000:.0f} k\n(marzo)',
                 xy=(i_min, v_min), xytext=(i_min+1.5, v_min*0.3),
                 fontsize=11, fontweight='bold', color=ROSSO,
                 arrowprops=dict(arrowstyle='->', color=ROSSO, lw=1.5))
    ax2.tick_params(axis='y', labelcolor=ROSSO)
    ax2.legend(loc='upper left', fontsize=11)

    fig.suptitle('Comportamento stagionale dei 4 sensori — cmos4 è anomalo',
                 fontsize=15, fontweight='bold', color=BLU, y=1.01)
    plt.tight_layout()
    return save_fig(fig, 'fig_cmos4_seasonal_drift.png')


def fig_pc1_corr_timescales():
    """r(PC1,T) e r(cmos1,T) a varie scale temporali."""
    scales = [('10min', '10 min'), ('1h', '1 ora'), ('1D', '1 giorno'),
              ('7D', '7 giorni'), ('30D', '30 giorni')]
    r_pc1_T, r_pc1_RH, r_c1_T = [], [], []
    for rule, _ in scales:
        d = df_raw[['PC1','cmos1','temperature','humidity']].resample(rule).mean().dropna()
        r_pc1_T.append(d['PC1'].corr(d['temperature']))
        r_pc1_RH.append(d['PC1'].corr(d['humidity']))
        r_c1_T.append(d['cmos1'].corr(d['temperature']))

    fig, ax = plt.subplots(figsize=(11, 5.5))
    x = np.arange(len(scales))
    w = 0.27
    b1 = ax.bar(x-w, np.abs(r_pc1_T), w, color=BLU, alpha=0.9,
                edgecolor='white', label='|r(PC₁, T)|')
    b2 = ax.bar(x, np.abs(r_pc1_RH), w, color=AZZURRO, alpha=0.9,
                edgecolor='white', label='|r(PC₁, RH)|')
    b3 = ax.bar(x+w, np.abs(r_c1_T), w, color=ROSSO, alpha=0.9,
                edgecolor='white', label='|r(cmos1, T)| [grezzo]')
    # Annotazioni
    for bars, vals in [(b1, r_pc1_T), (b2, r_pc1_RH), (b3, r_c1_T)]:
        for bar, v in zip(bars, vals):
            ax.annotate(f'{v:+.2f}', xy=(bar.get_x()+bar.get_width()/2, abs(v)),
                        xytext=(0, 3), textcoords='offset points',
                        ha='center', fontsize=9, color=GRIGIO)
    ax.set_xticks(x)
    ax.set_xticklabels([lbl for _, lbl in scales])
    ax.set_xlabel('Scala temporale di aggregazione')
    ax.set_ylabel('|Pearson r|')
    ax.set_title('Correlazione di PC₁ con T e RH al variare della scala temporale')
    ax.legend(loc='upper left', fontsize=10)
    ax.axhline(0.1, color=AZZURRO, linestyle=':', linewidth=1, alpha=0.5)
    ax.axhline(0.3, color=ROSSO, linestyle=':', linewidth=1, alpha=0.5)
    ax.text(0.98, 0.95,
            'PC₁ decorrelato a scala rapida (minuti-ore)\n'
            'Su scala mensile emerge la deriva stagionale (cmos4)',
            transform=ax.transAxes, ha='right', va='top', fontsize=10,
            bbox=dict(boxstyle='round,pad=0.5', facecolor=GRIG_LT,
                      edgecolor=GRIGIO, alpha=0.85))
    return save_fig(fig, 'fig_pc1_corr_timescales.png')


def fig_pc1_corr_monthly():
    """r(PC1,T) per singolo mese (scala oraria)."""
    df_h = df_raw[['PC1','temperature']].resample('1h').mean().dropna()
    df_h['month'] = df_h.index.month
    months_order = [3,4,5,6,7,8,9,10,11,12]
    labels = ['Mar','Apr','Mag','Giu','Lug','Ago','Set','Ott','Nov','Dic']
    rs = []
    for m in months_order:
        dm = df_h[df_h['month']==m]
        rs.append(dm['PC1'].corr(dm['temperature']))

    fig, ax = plt.subplots(figsize=(11, 5.5))
    colors = [ROSSO if abs(r) > 0.5 else AZZURRO if abs(r) > 0.2 else BLU
              for r in rs]
    xs = np.arange(len(months_order))
    bars = ax.bar(xs, rs, color=colors, alpha=0.9,
                  edgecolor='white', linewidth=1.5)
    for b, r in zip(bars, rs):
        ypos = r + (0.03 if r >= 0 else -0.06)
        ax.annotate(f'{r:+.2f}', xy=(b.get_x()+b.get_width()/2, ypos),
                    ha='center', fontsize=11, fontweight='bold',
                    color=GRIGIO)
    ax.axhline(0, color=GRIGIO, linewidth=0.8)
    ax.set_xticks(xs)
    ax.set_xticklabels(labels)
    ax.set_ylabel('r(PC₁, T) per mese')
    ax.set_title('Correlazione PC₁ vs T — per mese (scala oraria)',
                 pad=12)
    ax.set_ylim(-1.0, 0.5)
    # Evidenzia novembre
    nov_idx = labels.index('Nov')
    ax.annotate('Novembre:\nbrusca transizione\nstagionale di cmos4',
                xy=(nov_idx, rs[nov_idx]),
                xytext=(nov_idx-2.2, -0.9),
                fontsize=10, color=ROSSO,
                bbox=dict(boxstyle='round,pad=0.4', facecolor='white',
                          edgecolor=ROSSO, alpha=0.9),
                arrowprops=dict(arrowstyle='->', color=ROSSO, lw=1.5))
    ax.text(0.02, 0.02,
            'Il valore globale (r=−0.03) media mesi di segno opposto.\n'
            'La correlazione residua è una deriva lenta, non una risposta istantanea.',
            transform=ax.transAxes, ha='left', va='bottom', fontsize=10,
            bbox=dict(boxstyle='round,pad=0.5', facecolor=GRIG_LT,
                      edgecolor=GRIGIO, alpha=0.8))
    return save_fig(fig, 'fig_pc1_corr_monthly.png')


def fig_two_timescales():
    """Schema concettuale: due scale, due problemi, due soluzioni."""
    fig, ax = plt.subplots(figsize=(13, 5.8))
    ax.axis('off')
    ax.set_xlim(0, 13); ax.set_ylim(0, 6)

    # Titolo
    ax.text(6.5, 5.5, 'Due scale temporali — Due problemi — Due soluzioni',
            ha='center', fontsize=16, fontweight='bold', color=BLU)

    # Box sinistro: confounding veloce
    box1 = FancyBboxPatch((0.3, 1.2), 6, 3.5, boxstyle='round,pad=0.2',
                          facecolor=AZZ_LIGHT, edgecolor=BLU,
                          linewidth=2, alpha=0.35)
    ax.add_patch(box1)
    ax.text(3.3, 4.3, 'VELOCE  (minuti–ore)', ha='center',
            fontsize=14, fontweight='bold', color=BLU)
    ax.text(3.3, 3.7,
            'Ciclo T/RH giornaliero\n'
            'Tutti i MOX salgono insieme (common-mode)\n'
            'Sul grezzo: r(cmos1, T) alto intra-giornaliero',
            ha='center', fontsize=11, color=GRIGIO)
    # Freccia
    ax.annotate('', xy=(3.3, 2.2), xytext=(3.3, 3.0),
                arrowprops=dict(arrowstyle='->', color=BLU, lw=2))
    ax.text(3.3, 1.7, 'PCA → PC₁\n(common-mode su PC₂/PC₃)',
            ha='center', fontsize=12, fontweight='bold', color=BLU,
            bbox=dict(boxstyle='round,pad=0.4', facecolor='white',
                      edgecolor=BLU, alpha=0.9))

    # Box destro: confounding lento
    box2 = FancyBboxPatch((6.7, 1.2), 6, 3.5, boxstyle='round,pad=0.2',
                          facecolor='#FDE7E7', edgecolor=ROSSO,
                          linewidth=2, alpha=0.5)
    ax.add_patch(box2)
    ax.text(9.7, 4.3, 'LENTO  (settimane–mesi)', ha='center',
            fontsize=14, fontweight='bold', color=ROSSO)
    ax.text(9.7, 3.7,
            'Deriva stagionale di cmos4 (×59 mar→set)\n'
            'Si proietta su PC₁ via loading −0.34\n'
            'PC₁ settembre strutturalmente più basso',
            ha='center', fontsize=11, color=GRIGIO)
    ax.annotate('', xy=(9.7, 2.2), xytext=(9.7, 3.0),
                arrowprops=dict(arrowstyle='->', color=ROSSO, lw=2))
    ax.text(9.7, 1.7, 'Rolling LOD: μ_roll(t) adattivo\n(σ globale fisso)',
            ha='center', fontsize=12, fontweight='bold', color=ROSSO,
            bbox=dict(boxstyle='round,pad=0.4', facecolor='white',
                      edgecolor=ROSSO, alpha=0.9))

    # Footer
    ax.text(6.5, 0.5,
            'Pipeline completa: PCA + Rolling LOD = copertura di entrambe le scale del confounding',
            ha='center', fontsize=11, style='italic', color=GRIGIO)
    return save_fig(fig, 'fig_two_timescales.png')


def fig_cmos4_scales():
    """Tabella visuale: scale/CV dei sensori."""
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.axis('off')

    sensors = ['cmos1', 'cmos2', 'cmos3', 'cmos4']
    means = [3478, 3383, 152, 131032]
    sigmas = [491, 1708, 14, 132031]
    cvs = [sg/m*100 for sg, m in zip(sigmas, means)]
    ranges_approx = ['×1.5', '×5', '×2', '×59']

    cols = ['Sensore', 'Media', 'σ (std dev)', 'CV = σ/μ', 'Range mar→dic']
    n_rows = len(sensors) + 1
    n_cols = len(cols)

    widths = [0.17, 0.17, 0.20, 0.16, 0.20]
    x_starts = [sum(widths[:i]) for i in range(len(widths))]

    # Header
    for j, col in enumerate(cols):
        ax.add_patch(plt.Rectangle((x_starts[j]+0.05, 0.84),
                                     widths[j], 0.1, facecolor=BLU))
        ax.text(x_starts[j]+0.05+widths[j]/2, 0.89, col,
                ha='center', va='center', color='white',
                fontsize=13, fontweight='bold')

    # Rows
    for i, (s, m, sg, cv, rng) in enumerate(zip(sensors, means,
                                                  sigmas, cvs, ranges_approx)):
        y = 0.84 - 0.12 * (i+1)
        is_cmos4 = (s == 'cmos4')
        rcol = '#FDE7E7' if is_cmos4 else (GRIG_LT if i % 2 else 'white')
        for j in range(n_cols):
            ax.add_patch(plt.Rectangle((x_starts[j]+0.05, y),
                                         widths[j], 0.1,
                                         facecolor=rcol,
                                         edgecolor=GRIG_LT, linewidth=0.5))
        vals = [s, f'{m:,}', f'{sg:,}', f'{cv:.0f}%', rng]
        colors_cells = [BLU if j == 0 else (ROSSO if is_cmos4 else GRIGIO)
                        for j in range(n_cols)]
        bolds = [True, False, False, is_cmos4, is_cmos4]
        for j, (v, c, bld) in enumerate(zip(vals, colors_cells, bolds)):
            ax.text(x_starts[j]+0.05+widths[j]/2, y+0.05, v,
                    ha='center', va='center', color=c,
                    fontsize=12, fontweight='bold' if bld else 'normal')

    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.text(0.5, 0.97, 'Scale dei 4 sensori nel dataset (9 mesi)',
            ha='center', fontsize=15, fontweight='bold', color=BLU)
    ax.text(0.5, 0.02,
            'cmos4 ha CV del 100% e range stagionale ×59: non è un MOX come gli altri.',
            ha='center', fontsize=11, style='italic', color=GRIGIO)
    return save_fig(fig, 'fig_cmos4_scales.png')


def fig_pc1_seasonal_bias():
    """PC1 medio per mese: evidenzia bias stagionale."""
    dfm = df_raw[['PC1']].resample('1D').mean().dropna()
    dfm['month'] = dfm.index.month
    months_order = [3,4,5,6,7,8,9,10,11,12]
    labels = ['Mar','Apr','Mag','Giu','Lug','Ago','Set','Ott','Nov','Dic']
    means = [dfm[dfm['month']==m]['PC1'].mean() for m in months_order]
    # contributo stimato di cmos4: (-0.337) * z4_mese
    cmos4_monthly_mean = df_raw['cmos4'].resample('1D').mean().dropna()
    cmos4_monthly_mean = cmos4_monthly_mean.groupby(
        cmos4_monthly_mean.index.month).mean()
    z4_month = (cmos4_monthly_mean - mu_pca[3]) / sigma_pca[3]
    contrib = (-coeff[3, 0]) * z4_month

    fig, ax = plt.subplots(figsize=(11, 5.5))
    xs = np.arange(len(months_order))
    b1 = ax.bar(xs - 0.2, means, 0.4, color=BLU, alpha=0.9,
                edgecolor='white', label='PC₁ medio effettivo')
    b2 = ax.bar(xs + 0.2, [contrib.loc[m] for m in months_order],
                0.4, color=ROSSO, alpha=0.9,
                edgecolor='white', label='Contributo di cmos4 (−0.337·z₄)')
    ax.axhline(0, color=GRIGIO, linewidth=0.8)
    ax.axhline(lod_lo, color=ROSSO, linestyle='--', linewidth=1.5, alpha=0.6,
               label=f'LOD⁻ = {lod_lo:.2f}')
    ax.set_xticks(xs)
    ax.set_xticklabels(labels)
    ax.set_ylabel('PC₁ (unità standardizzate)')
    ax.set_title('Bias stagionale di PC₁ — dominato dal contributo di cmos4')
    ax.legend(loc='lower left', fontsize=10)
    # Annotazioni
    min_m_idx = np.argmin(means)
    ax.annotate(f'PC₁ depresso di\n{means[min_m_idx]:.2f} senza odori',
                xy=(min_m_idx, means[min_m_idx]),
                xytext=(min_m_idx+1.5, means[min_m_idx]-1.2),
                fontsize=10, color=BLU,
                bbox=dict(boxstyle='round,pad=0.4', facecolor='white',
                          edgecolor=BLU, alpha=0.9),
                arrowprops=dict(arrowstyle='->', color=BLU, lw=1.5))
    return save_fig(fig, 'fig_pc1_seasonal_bias.png')


# ───────────────────────────────────────────────────────────────────────────
# GENERA TUTTE LE FIGURE
# ───────────────────────────────────────────────────────────────────────────
print("\nGenerazione figure...")
fig_paths = {
    'loadings'          : fig_loadings(),
    'scree'             : fig_scree(),
    'baseline_cmp'      : fig_baseline_comparison(),
    'variance_decomp'   : fig_variance_decomposition(),
    'sigma_windows'     : fig_sigma_windows(),
    'rolling_events'    : fig_rolling_events(),
    'raw_events'        : fig_raw_event_rates(),
    'raw_temp_corr'     : fig_raw_temp_correlation(),
    'raw_false_pos'     : fig_raw_false_positives(),
    'k_sensitivity'     : fig_k_sensitivity(),
    'mox_decomp'        : fig_mox_decomposition(),
    'pipeline'          : fig_pipeline(),
    # nuove
    'cross_corr_bl'     : fig_cross_correlation_baseline(),
    'cmos4_seasonal'    : fig_cmos4_seasonal_drift(),
    'cmos4_scales'      : fig_cmos4_scales(),
    'pc1_corr_scales'   : fig_pc1_corr_timescales(),
    'pc1_corr_monthly'  : fig_pc1_corr_monthly(),
    'two_timescales'    : fig_two_timescales(),
    'pc1_seasonal_bias' : fig_pc1_seasonal_bias(),
}

# PNG riutilizzati
fig_paths['pc1_overview']    = EVT_DIR / 'PC1_overview.png'
fig_paths['pc1_vs_trh']      = EVT_DIR / 'PC1_vs_TRH_correlation.png'
fig_paths['rolling_overview']= ROL_DIR / 'PC1_fixed_vs_rolling.png'
fig_paths['signal_cmp']      = RVS_DIR / 'signal_comparison.png'
fig_paths['cmos4_vs_temp']   = FIGDIR / 'cmos4_vs_temperature.png'  # già generato

# ═══════════════════════════════════════════════════════════════════════════
# COSTRUZIONE .PPTX
# ═══════════════════════════════════════════════════════════════════════════
print("\nCostruzione presentazione...")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_filled_rect(slide, x, y, w, h, fill_rgb):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, *, size=18, bold=False,
                color=RGB_NERO, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullet_text(slide, x, y, w, h, bullets, *, size=16,
                    color=RGB_NERO, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f'• {b}'
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_header(slide, title):
    add_filled_rect(slide, Emu(0), Emu(0), SW, Inches(0.9), RGB_BLU)
    add_filled_rect(slide, Emu(0), Inches(0.9), SW, Inches(0.08), RGB_AZZURRO)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6),
                title, size=24, bold=True, color=RGB_BIANCO,
                anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page_num, total):
    add_filled_rect(slide, Emu(0), Inches(7.25), SW, Inches(0.25), RGB_GRIG_LT)
    add_textbox(slide, Inches(0.3), Inches(7.25), Inches(6), Inches(0.25),
                'Marco Calì — Progetto ELLONA', size=10, color=RGB_GRIGIO,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(10.5), Inches(7.25), Inches(2.5), Inches(0.25),
                f'{page_num} / {total}', size=10, color=RGB_GRIGIO,
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_image_centered(slide, img_path, *, top=Inches(1.3),
                       height=Inches(5.5)):
    img_path = str(img_path)
    from PIL import Image
    im = Image.open(img_path)
    img_w_px, img_h_px = im.size
    aspect = img_w_px / img_h_px
    height_emu = height
    width_emu  = Emu(int(height_emu * aspect))
    left = Emu((SW - width_emu) // 2)
    slide.shapes.add_picture(img_path, left, top, height=height_emu)


def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_filled_rect(s, Emu(0), Inches(4.5), SW, Inches(3), RGB_BLU)
    add_filled_rect(s, Emu(0), Inches(4.4), SW, Inches(0.1), RGB_AZZURRO)
    add_textbox(s, Inches(0.8), Inches(0.6), Inches(12), Inches(0.4),
                'Politecnico di Milano — Tesi di Laurea Magistrale',
                size=14, bold=True, color=RGB_BLU)
    add_textbox(s, Inches(0.8), Inches(1.7), Inches(12), Inches(1.5),
                'Rilevamento di anomalie olfattive',
                size=44, bold=True, color=RGB_BLU)
    add_textbox(s, Inches(0.8), Inches(2.8), Inches(12), Inches(1),
                'con electronic nose e analisi multivariata',
                size=30, bold=False, color=RGB_AZZURRO)
    add_textbox(s, Inches(0.8), Inches(3.7), Inches(12), Inches(0.6),
                'Pipeline PCA + Rolling LOD: due scale, due soluzioni',
                size=20, bold=False, color=RGB_GRIGIO)
    add_textbox(s, Inches(0.8), Inches(5), Inches(10), Inches(0.5),
                'Progetto ELLONA · Monitoraggio odori · PoliMi — IREN',
                size=15, bold=True, color=RGB_BIANCO)
    add_textbox(s, Inches(0.8), Inches(5.7), Inches(10), Inches(0.5),
                'Marco Calì',
                size=22, bold=True, color=RGB_BIANCO)
    add_textbox(s, Inches(0.8), Inches(6.3), Inches(10), Inches(0.5),
                'Aprile 2026 — Avanzamento lavori',
                size=16, bold=False, color=RGB_AZZ_LT)
    return s


def slide_agenda():
    s = prs.slides.add_slide(BLANK)
    add_header(s, 'Indice')
    items = [
        ('1.', 'Contesto del progetto',         'ELLONA, impianto, dispositivo, dataset'),
        ('2.', 'Il problema della detection',   'MOX, confounders, dimensionalità'),
        ('3.', 'La PCA come prima difesa',      'Principio, loading empirici, PC₁'),
        ('4.', 'La scoperta di cmos4',          'Anomalia stagionale, cross-correlazione'),
        ('5.', 'Due scale temporali',           'Confounding veloce vs lento'),
        ('6.', 'Costruzione del LOD su PC₁',    'Baseline IQR, scelta di k'),
        ('7.', 'Il Rolling LOD',                'Deriva stagionale, decomposizione σ, Opzione B'),
        ('8.', 'Confronto grezzo vs PC₁',       'Dimostrazione empirica dei falsi positivi'),
        ('9.', 'Conclusioni e prossimi passi',  'Limitazioni aperte, lavoro futuro'),
    ]
    y = Inches(1.3)
    for num, title, desc in items:
        add_textbox(s, Inches(0.8), y, Inches(0.8), Inches(0.58),
                    num, size=22, bold=True, color=RGB_AZZURRO)
        add_textbox(s, Inches(1.6), y, Inches(5), Inches(0.58),
                    title, size=18, bold=True, color=RGB_BLU,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, Inches(6.8), y, Inches(6.2), Inches(0.58),
                    desc, size=13, color=RGB_GRIGIO, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.63)
    return s


def slide_section(number, title, subtitle):
    s = prs.slides.add_slide(BLANK)
    add_filled_rect(s, Emu(0), Emu(0), SW, SH, RGB_BLU)
    add_filled_rect(s, Emu(0), Inches(3.2), SW, Inches(0.1), RGB_AZZ_LT)
    add_textbox(s, Inches(1), Inches(2), Inches(11), Inches(1),
                f'Parte {number}', size=20, bold=True, color=RGB_AZZ_LT)
    add_textbox(s, Inches(1), Inches(3.5), Inches(11), Inches(1.5),
                title, size=40, bold=True, color=RGB_BIANCO)
    add_textbox(s, Inches(1), Inches(5), Inches(11), Inches(1),
                subtitle, size=18, bold=False, color=RGB_AZZ_LT)
    return s


def slide_bullets(title, bullets, *, sub=None):
    s = prs.slides.add_slide(BLANK)
    add_header(s, title)
    y0 = Inches(1.4)
    if sub:
        add_textbox(s, Inches(0.8), y0, Inches(11.8), Inches(0.6),
                    sub, size=16, bold=False, color=RGB_AZZURRO)
        y0 += Inches(0.7)
    add_bullet_text(s, Inches(0.8), y0, Inches(11.8), Inches(5.5),
                    bullets, size=18)
    return s


def slide_bullets_image(title, bullets, img_path, *, sub=None,
                        img_width=Inches(6.5)):
    s = prs.slides.add_slide(BLANK)
    add_header(s, title)
    y0 = Inches(1.3)
    if sub:
        add_textbox(s, Inches(0.5), y0, Inches(12.3), Inches(0.5),
                    sub, size=15, color=RGB_AZZURRO)
        y0 += Inches(0.55)
    add_bullet_text(s, Inches(0.5), y0, Inches(5.6), Inches(5.5),
                    bullets, size=15)
    try:
        from PIL import Image
        im = Image.open(str(img_path))
        aspect = im.size[0] / im.size[1]
        h = Inches(5)
        w = Emu(int(h * aspect))
        if w > img_width:
            w = img_width
            h = Emu(int(w / aspect))
        left = Inches(13.333) - w - Inches(0.4)
        s.shapes.add_picture(str(img_path), left, y0 + Inches(0.2),
                              width=w, height=h)
    except Exception as e:
        print(f"  Warning: image {img_path} not added: {e}")
    return s


def slide_image_full(title, img_path, *, sub=None):
    s = prs.slides.add_slide(BLANK)
    add_header(s, title)
    y0 = Inches(1.3)
    if sub:
        add_textbox(s, Inches(0.5), y0, Inches(12.3), Inches(0.5),
                    sub, size=15, color=RGB_AZZURRO, align=PP_ALIGN.CENTER)
        y0 += Inches(0.55)
    add_image_centered(s, img_path, top=y0 + Inches(0.1),
                       height=Inches(5.3))
    return s


def slide_table(title, headers, rows, *, sub=None, col_widths=None):
    s = prs.slides.add_slide(BLANK)
    add_header(s, title)
    y0 = Inches(1.3)
    if sub:
        add_textbox(s, Inches(0.5), y0, Inches(12.3), Inches(0.5),
                    sub, size=15, color=RGB_AZZURRO)
        y0 += Inches(0.55)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = Inches(12.3)
    if col_widths is None:
        col_widths = [total_w // n_cols] * n_cols
    table_h = Inches(0.5) + Inches(0.45) * len(rows)
    tbl = s.shapes.add_table(n_rows, n_cols, Inches(0.5), y0 + Inches(0.1),
                              total_w, table_h).table
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGB_BLU
        cell.text = ''
        tf = cell.text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = RGB_BIANCO; r.font.name = 'Calibri'
    for i, row in enumerate(rows):
        row_color = RGB_GRIG_LT if i % 2 == 1 else RGB_BIANCO
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = row_color
            cell.text = ''
            tf = cell.text_frame
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(13); r.font.name = 'Calibri'
            r.font.color.rgb = RGB_GRIGIO
            if j == 0:
                r.font.bold = True; r.font.color.rgb = RGB_BLU
    return s


def slide_end():
    s = prs.slides.add_slide(BLANK)
    add_filled_rect(s, Emu(0), Emu(0), SW, SH, RGB_BLU)
    add_textbox(s, Inches(0), Inches(2.8), Inches(13.333), Inches(1.2),
                'Grazie per l\'attenzione',
                size=50, bold=True, color=RGB_BIANCO,
                align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0), Inches(4.2), Inches(13.333), Inches(0.7),
                'Domande e discussione',
                size=24, color=RGB_AZZ_LT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0), Inches(5.5), Inches(13.333), Inches(0.5),
                'Marco Calì   ·   Progetto ELLONA   ·   PoliMi — IREN',
                size=14, color=RGB_AZZ_LT, align=PP_ALIGN.CENTER)
    return s


# ───────────────────────────────────────────────────────────────────────────
# COSTRUZIONE SEQUENZIALE SLIDE
# ───────────────────────────────────────────────────────────────────────────
slides = []

# Title + Indice
slides.append(slide_title())
slides.append(slide_agenda())

# ═══ PARTE 1 — CONTESTO ═══
slides.append(slide_section('1', 'Contesto del progetto',
    'ELLONA, l\'impianto, il dispositivo, il dataset'))

slides.append(slide_bullets(
    'Il progetto ELLONA',
    [
        'Collaborazione Politecnico di Milano — IREN',
        'Obiettivo: automazione del rilevamento di eventi olfattivi presso impianti di trattamento',
        'Approccio: electronic nose (array di 4 sensori MOX) + analisi dati in tempo reale',
        'Problematica: odori emessi dagli impianti hanno forte impatto sulla popolazione limitrofa',
        'Stato dell\'arte: segnalazioni manuali + ispezioni olfattometriche — servono strumenti oggettivi e continui',
    ], sub='Monitoraggio continuo e non invasivo della qualità dell\'aria'))

slides.append(slide_bullets(
    'L\'impianto di monitoraggio',
    [
        'Impianto di trattamento acque reflue, Nord Italia',
        'Sorgenti principali di odori: metabolismo batterico nei fanghi, digestione anaerobica',
        'Picco di emissioni in estate: temperatura e attività biologica favoriscono il rilascio di composti volatili (H₂S, NH₃, VOC)',
        'Dispositivo ELLONA installato in posizione strategica — esposizione diretta ai pennacchi odorigeni',
    ]))

slides.append(slide_bullets(
    'Il dispositivo ELLONA',
    [
        '4 sensori MOX (metal oxide): cmos1, cmos2, cmos3, cmos4 — tecnologie e sensibilità chimiche diverse',
        'Sensori ambientali di riferimento: temperatura (T), umidità relativa (RH)',
        'Sensori specifici (NH₃, H₂S) — non usati per baseline PCA (specifici a singolo analita)',
        'Acquisizione sincrona a frequenza 10 s',
        'Deploy continuo in situ da marzo 2025',
    ]))

slides.append(slide_bullets(
    'Il dataset di lavoro',
    [
        'Periodo: marzo → dicembre 2025 (9 mesi continui)',
        'Campioni totali: ~2.069.000 @ 10 s',
        'Pre-processing: parsing CSV mensili, concatenazione, rimozione duplicati e NaN',
        'Variabili conservate: cmos1–4, T, RH, NH₃, H₂S',
        'Copertura temporale ≈ 100% — pochi gap di acquisizione',
    ], sub='2 milioni di campioni su 9 mesi: base statisticamente solida'))

# ═══ PARTE 2 — IL PROBLEMA DELLA DETECTION ═══
slides.append(slide_section('2', 'Il problema della detection',
    'Perché una soglia sul segnale grezzo non funziona'))

slides.append(slide_bullets(
    'Il Limit of Detection (LOD): cos\'è',
    [
        'LOD = soglia statistica che separa "normalità" da "evento"',
        'Formulazione classica: LOD = μ ± k·σ  dove μ, σ sono media e deviazione standard del baseline',
        'k è il moltiplicatore: determina la specificità (tipicamente k = 3 → ~99.7% su distribuzione gaussiana)',
        'Un evento è dichiarato quando il segnale esce dalla banda [μ − kσ,  μ + kσ]',
        'Problema chiave: su QUALE SEGNALE applicare il LOD?',
    ], sub='Il LOD è un classificatore binario; la qualità dipende dal segnale di ingresso'))

slides.append(slide_image_full(
    'Il MOX come trasduttore non selettivo',
    fig_paths['mox_decomp'],
    sub='Il segnale MOX è una miscela di contributi inseparabili — soglia diretta = falsi positivi'))

slides.append(slide_bullets(
    'Problema #1 — La dimensionalità',
    [
        'ELLONA ha 4 sensori MOX → 4 segnali grezzi, 4 possibili soglie',
        'Come si combinano in un\'unica decisione?',
        '    — AND: alta specificità, bassa sensibilità (basta un sensore non rispondente)',
        '    — OR: alta sensibilità, molti falsi positivi',
        '    — Media: perde l\'informazione di correlazione tra sensori',
        'Serve un approccio che sfrutti la STRUTTURA DI CORRELAZIONE dei sensori',
    ], sub='I sensori di un e-nose sono correlati: ignorarlo = perdere informazione'))

slides.append(slide_bullets_image(
    'Problema #2 — I confounders ambientali',
    [
        'Temperatura stagionale: 5–35 °C',
        'Umidità relativa: 40–95 %',
        'Entrambe producono variazioni di resistenza MOX dello stesso ordine della risposta agli odori',
        'Drift lento di invecchiamento: la risposta cambia nel tempo',
        'Un LOD sul grezzo è dominato da questi effetti → falsi positivi stagionali',
    ],
    fig_paths['baseline_cmp'],
    sub='Il baseline globale sul grezzo produce 13% eventi: dominato dalla deriva stagionale'))

# ═══ PARTE 3 — LA PCA COME PRIMA DIFESA ═══
slides.append(slide_section('3', 'La PCA come prima difesa',
    'Rimuovere il confounding rapido di T/RH'))

slides.append(slide_bullets(
    'La PCA: cosa fa',
    [
        'Principal Component Analysis: trasformazione lineare che trova nuovi assi lungo i quali la varianza dei dati è massima',
        'Matematicamente: le PC sono gli autovettori della matrice di covarianza, ordinati per autovalore decrescente',
        'Proprietà chiave: i nuovi assi sono ortogonali (decorrelati) tra loro',
        'Nel nostro caso: 4 sensori MOX → 4 componenti principali (PC₁, PC₂, PC₃, PC₄)',
        'PC₁ = direzione di massima varianza nei dati — NON necessariamente la "direzione dell\'odore"',
        'Cosa cattura PC₁ nel nostro caso? → va verificato empiricamente',
    ], sub='La PCA trova direzioni di massima varianza — non sa a priori cosa è odore'))

slides.append(slide_image_full(
    'Pipeline ELLONA — visione d\'insieme',
    fig_paths['pipeline']))

slides.append(slide_bullets_image(
    'Loading plot: risultato empirico della PCA',
    [
        'I loading di PC₁ sono i coefficienti con cui ciascun sensore contribuisce a PC₁',
        'cmos1: +0.528     cmos2: +0.669',
        'cmos3: +0.401     cmos4: −0.336',
        'Il segno negativo di cmos4 è inatteso — lo abbiamo scoperto, non progettato',
        'Interpretazione immediata: cmos4 si muove "al contrario" degli altri nel dataset — perché?',
        'Questa domanda è il punto di partenza della Parte 4',
    ],
    fig_paths['loadings'],
    sub='PC₁ cattura il 44% della varianza — ma cosa significa quel loading negativo?'))

slides.append(slide_bullets_image(
    'Varianza spiegata dalle componenti',
    [
        'PC₁: 44% — PC₂: 26% — PC₃: 21% — PC₄: 9%',
        'PC₁ + PC₂ ≈ 70%: le prime due componenti catturano la parte dominante',
        'Per il LOD operativo si usa solo PC₁',
        'PC₂ e PC₃ utili per diagnosi e possibile fingerprinting degli odori (lavoro futuro)',
    ],
    fig_paths['scree']))

# ═══ PARTE 4 — LA SCOPERTA DI CMOS4 ═══
slides.append(slide_section('4', 'La scoperta di cmos4',
    'Il loading negativo non è ciò che sembra'))

slides.append(slide_bullets(
    'Il primo sospetto: "cmos4 risponde al contrario agli odori"',
    [
        'Interpretazione iniziale: loading negativo → cmos4 sale quando gli altri scendono',
        'Coerente con l\'idea di "asse differenziale" della PCA',
        'Ma osservazione empirica in campo: durante l\'esposizione a un odore TUTTI e 4 i sensori scendono insieme',
        'Inconsistenza: se tutti scendono, perché cmos4 ha segno opposto?',
        'Serve guardare i dati grezzi per capire cosa cattura davvero PC₁',
    ], sub='Domanda aperta: da dove viene il loading negativo?'))

slides.append(slide_image_full(
    'Evidenza #1 — Correlazione tra sensori nel baseline',
    fig_paths['cross_corr_bl'],
    sub='cmos4 è anti-correlato con cmos2 (r=−0.33) nel baseline: NON è un effetto di breve termine'))

slides.append(slide_image_full(
    'Evidenza #2 — cmos4 ha una deriva stagionale anomala',
    fig_paths['cmos4_seasonal'],
    sub='cmos1 ±20% · cmos2 ×5 · cmos3 ±25% · cmos4 ×59 — ordine di grandezza diverso'))

slides.append(slide_image_full(
    'Le scale dei 4 sensori',
    fig_paths['cmos4_scales'],
    sub='cmos4 ha CV = 100% (σ ≈ media): comportamento incompatibile con un normale sensore MOX'))

slides.append(slide_bullets(
    'Interpretazione: cmos4 è un sensore diverso',
    [
        'Il loading negativo NON riflette una risposta opposta agli odori',
        'Riflette l\'anti-correlazione stagionale tra cmos4 (estate: ×59) e cmos2 (estate: −65%)',
        'Durante un evento odorigeno tutti i 4 sensori scendono — il contributo di cmos4 ATTENUA parzialmente PC₁ (non lo amplifica)',
        'Ipotesi fisica: cmos4 è un sensore per gas ossidanti (NOₓ, O₃) — in estate più UV → più O₃ → R cmos4 sale drammaticamente',
        'Domanda aperta per la tesi: caratterizzare meglio il ruolo di cmos4 nell\'array',
    ], sub='Questa è una scoperta, non un difetto — ma cambia l\'interpretazione della PCA'))

# ═══ PARTE 5 — LE DUE SCALE TEMPORALI ═══
slides.append(slide_section('5', 'Due scale temporali',
    'Confounding veloce e confounding lento'))

slides.append(slide_image_full(
    'r(PC₁, T) cambia con la scala temporale',
    fig_paths['pc1_corr_scales'],
    sub='PC₁ è decorrelato a scala rapida. Su scala mensile emerge la correlazione con RH (cmos4)'))

slides.append(slide_image_full(
    'r(PC₁, T) per singolo mese (scala oraria)',
    fig_paths['pc1_corr_monthly'],
    sub='I valori mensili oscillano da −0.79 a +0.30: il bias stagionale di cmos4 domina'))

slides.append(slide_image_full(
    'Il bias stagionale di PC₁ — contributo di cmos4',
    fig_paths['pc1_seasonal_bias'],
    sub='Il contributo (−0.337·z₄) spiega la quasi totalità della variazione mensile di PC₁'))

slides.append(slide_image_full(
    'Due scale, due problemi, due soluzioni',
    fig_paths['two_timescales'],
    sub='Questa è la struttura metodologica della tesi'))

# ═══ PARTE 6 — LOD SU PC1 ═══
slides.append(slide_section('6', 'Costruzione del LOD su PC₁',
    'Baseline IQR, scelta di k, applicazione'))

slides.append(slide_bullets(
    'La selezione del baseline con filtro IQR',
    [
        'Problema: non tutti i campioni rappresentano condizioni normali — molti contengono già eventi',
        'Includerli nel baseline distorce μ e σ → LOD meno sensibile',
        'Soluzione: filtro IQR [P25, P75] applicato a TUTTI e 4 i sensori MOX, settimanalmente',
        'Un campione è "baseline" se ciascuno dei 4 sensori cade nella propria banda IQR (intersezione)',
        'Robusto a outlier: i percentili non sono affetti dai picchi di eventi',
    ], sub='~340.000 campioni qualificano come baseline (16.5%) — base per μ e σ'))

slides.append(slide_image_full(
    'Confronto dei baseline modes',
    fig_paths['baseline_cmp'],
    sub='Weekly è la scelta ottimale: cattura deriva lenta senza adattarsi agli eventi'))

slides.append(slide_bullets(
    'La soglia LOD: μ ± k·σ',
    [
        'Su PC₁ del baseline: μ_BL ≈ 0 (per costruzione dopo z-score), σ_BL = 1.329',
        'LOD⁻ = μ_BL − k·σ_BL = −3.988    (k = 3)',
        'LOD⁺ = μ_BL + k·σ_BL = +3.988',
        'Solo LOD⁻ è operativo: la resistenza MOX cala con i gas riducenti (odori tipici)',
        'Evento ⟺ PC₁(t) < LOD⁻',
    ], sub='Risultato: 3.67% dei 2M campioni qualifica come evento'))

slides.append(slide_image_full(
    'Applicazione: PC₁(t) con LOD su 9 mesi',
    fig_paths['pc1_overview'],
    sub='Picchi concentrati in agosto–settembre — coerenti con il picco di attività biologica estiva'))

# ═══ PARTE 7 — ROLLING LOD ═══
slides.append(slide_section('7', 'Il Rolling LOD',
    'Compensare la deriva stagionale di cmos4'))

slides.append(slide_bullets(
    'Perché il LOD fisso non basta (a lungo termine)',
    [
        'μ_BL, σ_BL calcolati una volta su tutto il dataset (assunzione di stazionarietà)',
        'Ma PC₁ ha una deriva stagionale: il contributo di cmos4 sposta μ_PC₁ nel tempo',
        'Se a settembre il PC₁ medio è −2.8 (senza odori!), la soglia LOD⁻ = −3.99 resta calibrata su un baseline diverso',
        'Su 9 mesi l\'effetto è contenuto, ma su 2+ anni (deployment operativo) diventa critico',
        'Serve un LOD che segua la deriva senza perdere calibrazione statistica',
    ]))

slides.append(slide_bullets(
    'Primo tentativo: LOD rolling con σ locale (FALLITO)',
    [
        'Idea: ricalcolare μ(t) e σ(t) ogni giorno su una finestra causale di 7 giorni',
        'Formula: LOD(t) = μ_roll(t) ± k · σ_roll(t)',
        'Risultato catastrofico: 23.04% eventi (vs 3.67% del LOD fisso)',
        'Diagnosi: σ_roll = 0.44 ≈ 3× più piccolo di σ_global = 1.33 → banda 3× più stretta → tutto diventa "evento"',
        'Il problema è strutturale, non un bug',
    ], sub='Perché σ_roll è così diverso da σ_global? Domanda centrale'))

slides.append(slide_image_full(
    'σ_roll al variare della finestra (7d / 14d / 30d)',
    fig_paths['sigma_windows'],
    sub='σ_roll non converge a σ_global nemmeno a 30 giorni (solo 41%)'))

slides.append(slide_bullets_image(
    'Decomposizione della varianza',
    [
        'σ²_global ≈ σ²_rumore + σ²_stagionale',
        'σ²_rumore ≈ 0.44² = 0.19 (rumore di breve termine)',
        'σ²_stagionale ≈ 1.33² − 0.19 = 1.58 (scala mesi)',
        '→ l\'89% di σ_global è variabilità stagionale',
        'Finestre 7/14/30 giorni vedono la stagionalità come offset costante (catturata da μ_roll)',
        'σ_roll ha definizione solo sul rumore veloce — non converge',
    ],
    fig_paths['variance_decomp'],
    sub='La stagionalità vive su scale > 30 giorni'))

slides.append(slide_bullets(
    'Analogia: la temperatura di Milano',
    [
        'Misuri σ della T di Milano su 1 settimana di luglio: 28–34°C → σ ≈ 2°C',
        'Misuri σ sull\'intero anno: 2–36°C → σ ≈ 10°C',
        'Stessa città, stesso termometro — σ diverso perché la finestra vede variabilità diversa',
        'In PC₁: σ_roll coglie il "rumore giornaliero" (2°C), σ_global coglie l\'intera variabilità (10°C)',
        'Per una soglia valida in ogni stagione serve il σ annuale, non quello settimanale',
    ], sub='La scelta di σ è una scelta di scala temporale'))

slides.append(slide_bullets(
    'Soluzione: Opzione B — μ rolling, σ globale',
    [
        'LOD(t) = μ_roll(t) ± k · σ_global',
        'μ_roll(t): aggiornato ogni giorno (ultima settimana) → segue la deriva stagionale di cmos4',
        'σ_global: fisso → calibrato sull\'intera variabilità stagionale',
        'Concetto analogo ai grafici SPC Moving Average: centro adattivo, limiti fissi',
        'Motivazione fisica: σ rappresenta l\'intera variabilità reale (89% stagionale), non il rumore di breve termine',
    ], sub='Rolling per correggere la deriva, fisso per preservare la calibrazione'))

slides.append(slide_image_full(
    'Risultati: LOD fisso vs rolling',
    fig_paths['rolling_events'],
    sub='Concordanza 98.5–98.9%: su 9 mesi la deriva di μ è ancora contenuta'))

slides.append(slide_image_full(
    'Sensibilità al parametro k',
    fig_paths['k_sensitivity'],
    sub='A k=3 il LOD fisso è già solido. A k<1.5 il rolling fa la differenza'))

# ═══ PARTE 8 — CONFRONTO GREZZO VS PC1 ═══
slides.append(slide_section('8', 'Confronto empirico: grezzo vs PC₁',
    'La dimostrazione che la PCA è necessaria'))

slides.append(slide_bullets(
    'La domanda fondamentale',
    [
        'Tutto il ragionamento è stato costruito su PC₁ — ma PERCHÉ non applicare il LOD direttamente sul grezzo?',
        'Argomento teorico: il grezzo contiene confounders T/RH',
        'Argomento empirico: servono numeri sui dati reali',
        'Setup: costruiamo un LOD su ciascun segnale grezzo con la stessa metodologia (baseline IQR, k=3)',
    ], sub='Script ELLONA_12_raw_vs_pca_lod.m: confronto sistematico'))

slides.append(slide_bullets(
    'Setup del confronto',
    [
        'Baseline mask comune: intersezione IQR [P25, P75] su tutti i 4 MOX (11.8% dei campioni)',
        'Stesso k = 3 per tutti i metodi',
        '6 metodi confrontati:',
        '    — cmos1, cmos2, cmos3, cmos4 (LOD diretto sul grezzo)',
        '    — avg_z (media z-normalizzata dei 4 MOX: combinazione naïve)',
        '    — PC₁ (combinazione PCA ottimale, riferimento)',
        'Tre metriche: tasso eventi, correlazione con temperatura, overlap con PC₁',
    ]))

slides.append(slide_image_full(
    'Risultato #1: tasso di eventi per metodo',
    fig_paths['raw_events'],
    sub='Inconsistenza tra sensori grezzi (0.3%–14%) sullo stesso fenomeno'))

slides.append(slide_image_full(
    'Risultato #2: correlazione settimanale con T',
    fig_paths['raw_temp_corr'],
    sub='cmos1 ha r=+0.508 settimanale; PC₁ ha r=+0.033 — la PCA rimuove la deriva veloce'))

slides.append(slide_image_full(
    'Risultato #3: i falsi positivi del grezzo',
    fig_paths['raw_false_pos'],
    sub='Il 97% degli eventi di cmos1 NON è confermato da PC₁ — quasi interamente FP stagionali'))

slides.append(slide_table(
    'Riepilogo quantitativo: grezzo vs PC₁',
    ['Metodo', '% eventi', 'r con T', 'Overlap PC₁', '% FP'],
    [
        ['cmos1',  '12.62%',  '+0.508',  '0.31%',   '97.6%'],
        ['cmos2',  '1.86%',   '−0.205',  '0.26%',   '86.0%'],
        ['cmos3',  '14.15%',  '−0.313',  '0.38%',   '97.3%'],
        ['cmos4',  '0.32%',   '−0.366',  '0.00%',   '100%'],
        ['avg_z',  '12.02%',  '−0.101',  '0.30%',   '97.6%'],
        ['PC₁',    '3.67%',   '+0.033',  '3.67%',   '—'],
    ],
    sub='LOD diretto sul grezzo = quasi esclusivamente falsi positivi stagionali'))

slides.append(slide_bullets(
    'Interpretazione: tre prove convergenti',
    [
        '1. Inconsistenza tra sensori — tassi da 0.3% a 14% sullo stesso fenomeno: non esiste una soglia "giusta" sul grezzo',
        '2. Correlazione con T — cmos1 r=+0.51, PC₁ r=+0.03: la PCA rimuove il confounding veloce',
        '3. Falsi positivi — 97% degli eventi cmos1 non confermati da PC₁: dominanza del rumore stagionale',
        'Le tre evidenze convergono: la PCA non è un artifizio ma una necessità fisica',
        'La deriva stagionale di cmos4 (scala lenta) è un problema separato, gestito dal Rolling LOD',
    ], sub='Il nucleo della tesi: due stadi di filtraggio per due scale temporali'))

# ═══ PARTE 9 — CONCLUSIONI ═══
slides.append(slide_section('9', 'Conclusioni e prossimi passi',
    'Risultati, limitazioni aperte, lavoro futuro'))

slides.append(slide_image_full(
    'Sintesi visiva: la pipeline completa',
    fig_paths['pipeline']))

slides.append(slide_table(
    'Numeri chiave della ricerca',
    ['Affermazione', 'Valore', 'Fonte'],
    [
        ['Campioni analizzati',                       '2.069.099',           'ELLONA_07'],
        ['Baseline IQR (weekly)',                     '341.820 (16.5%)',     'ELLONA_08'],
        ['σ di PC₁ sul baseline',                     '1.329',               'ELLONA_08'],
        ['Soglia LOD⁻ (k=3)',                         '−3.988',              'ELLONA_08'],
        ['Tasso eventi LOD su PC₁',                   '3.67%',               'ELLONA_08'],
        ['Tasso eventi LOD su cmos1 grezzo',          '12.62%',              'ELLONA_12'],
        ['FP cmos1 (non conf. da PC₁)',               '97.6%',               'ELLONA_12'],
        ['r(cmos1, T) settimanale',                   '+0.508',              'ELLONA_12'],
        ['r(PC₁, T) a 10 min',                        '−0.032',              'analisi post'],
        ['r(PC₁, RH) mensile',                        '−0.557',              'analisi post'],
        ['Deriva stagionale cmos4 (mar→set)',         '×59 (20k → 1.2M)',    'analisi post'],
        ['Varianza stagionale in σ_global',           '89%',                 'ELLONA_11'],
        ['Concordanza LOD fisso vs rolling (k=3)',    '98.5%',               'ELLONA_11'],
    ],
    col_widths=[Inches(6.5), Inches(3.2), Inches(2.6)]))

slides.append(slide_bullets(
    'Cosa abbiamo dimostrato',
    [
        'Il segnale MOX grezzo è strutturalmente inadatto per un LOD diretto (97% FP)',
        'La PCA produce un segnale (PC₁) decorrelato dai confounders ambientali VELOCI (minuti-ore)',
        'La selezione del baseline via filtro IQR è robusta e garantisce calibrazione su condizioni normali',
        'Su scala stagionale PC₁ conserva una deriva — dovuta principalmente a cmos4 — che il Rolling LOD compensa',
        'Due stadi di filtraggio per due scale temporali: PCA (veloce) + Rolling LOD (lento)',
    ]))

slides.append(slide_bullets(
    'Limitazioni aperte',
    [
        'cmos4 resta un sensore anomalo: la sua natura fisica non è stata caratterizzata (ipotesi: gas ossidanti)',
        'Il contributo di cmos4 a PC₁ è negativo ma non trascurabile (loading −0.336)',
        'PC₁ non è "decorrelato da T/RH" in assoluto: su scala mensile r(PC₁, RH) = −0.557',
        'Manca un ground truth: non abbiamo ancora validato gli eventi contro segnalazioni olfattometriche',
        'Deployment esteso (2+ anni) è necessario per valutare l\'effettivo vantaggio del Rolling LOD',
    ], sub='Onestà metodologica: ogni affermazione ha una scala temporale di validità'))

slides.append(slide_bullets(
    'Prossimi passi',
    [
        'Caratterizzazione di cmos4: identificare tipo di sensore, risposta a NOₓ/O₃, ruolo nell\'array',
        'Valutare se escludere cmos4 o riprogettare la PCA sui soli 3 sensori concordi',
        'Validazione con ground truth: correlazione eventi con segnalazioni olfattometriche o campionamenti',
        'Classificazione degli odori: ogni evento ha un pattern in PC₂/PC₃ → possibile fingerprinting',
        'GUI operativa e alerting in tempo reale (ELLONA_10 in sviluppo)',
    ]))

# Chiusura
slides.append(slide_end())

# Footer su tutte tranne prima e ultima
total = len(slides)
for i, s in enumerate(slides):
    if 0 < i < total - 1:
        add_footer(s, i + 1, total)

pptx_path = OUTDIR / 'thesis_progress.pptx'
prs.save(str(pptx_path))

print(f"\n✓ Salvato: {pptx_path}")
print(f"  Numero slide: {total}")
n_gen = sum(1 for p in fig_paths.values() if 'thesis_pptx/figures' in str(p))
print(f"  Figure generate: {n_gen}")
print(f"  Figure riusate:  {len(fig_paths) - n_gen}")
