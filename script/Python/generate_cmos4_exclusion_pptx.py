"""
generate_cmos4_exclusion_pptx.py
─────────────────────────────────
PowerPoint focalizzato sulla motivazione empirica dell'esclusione di cmos4.

Struttura (10 slide):
  1. Titolo
  2. Il problema: deriva ×59 e loading negativo
  3. Ipotesi umidità: AH non spiega la deriva (R²≈0)
  4. Ipotesi T e RH separati: nessuna variabile spiega cmos4
  5. Profili stagionali: il timing di cmos4 non batte né T né RH
  6. Conseguenza sulla PCA: loading plot 3 vs 4 sensori
  7. Bias stagionale di PC1: si riduce ma non sparisce
  8. Impatto sulla detection: −67% falsi positivi
  9. Stazionarietà della baseline: σ_roll/σ_global → 1
 10. Conclusione e decisione

Marco Calì — PoliMi, Aprile 2026
"""

from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── PATHS ─────────────────────────────────────────────────────────────────
BASE   = Path("/Users/marcocali/Desktop/IREN")
OUTDIR = BASE / "output" / "cmos4_exclusion_pptx"
OUTDIR.mkdir(parents=True, exist_ok=True)

FIG_13 = BASE / "output" / "humidity_compensation"
FIG_14 = BASE / "output" / "env_correlation"
FIG_15 = BASE / "output" / "pca_3vs4"
FIG_ED = BASE / "output" / "event_detection"   # figure originali

# ── PALETTE POLIMI ─────────────────────────────────────────────────────────
BLU       = "#003866"
AZZURRO   = "#1464A0"
AZZ_LIGHT = "#6BA3D6"
GRIGIO    = "#58595B"
GRIG_LT   = "#D0D2D3"
ROSSO     = "#B5394E"
ARANCIO   = "#E8A33D"
VERDE     = "#4A9B6E"

RGB_BLU    = RGBColor(0x00, 0x38, 0x66)
RGB_AZZ    = RGBColor(0x14, 0x64, 0xA0)
RGB_AZZ_LT = RGBColor(0x6B, 0xA3, 0xD6)
RGB_GRIGIO = RGBColor(0x58, 0x59, 0x5B)
RGB_GRIG_LT= RGBColor(0xD0, 0xD2, 0xD3)
RGB_ROSSO  = RGBColor(0xB5, 0x39, 0x4E)
RGB_ARANCIO= RGBColor(0xE8, 0xA3, 0x3D)
RGB_VERDE  = RGBColor(0x4A, 0x9B, 0x6E)
RGB_BIANCO = RGBColor(0xFF, 0xFF, 0xFF)
RGB_NERO   = RGBColor(0x1A, 0x1A, 0x1A)

# ── PRESENTAZIONE ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK  = prs.slide_layouts[6]


# ── HELPER FUNCTIONS ──────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, rgb):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = rgb
    s.line.fill.background()
    return s


def add_text(slide, x, y, w, h, text, *, size=18, bold=False,
             color=RGB_NERO, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top  = tf.margin_bottom = Emu(0)
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name   = 'Calibri'
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=16, color=RGB_NERO):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = f'• {item}'
        run.font.name  = 'Calibri'
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None):
    add_rect(slide, Emu(0), Emu(0), SW, Inches(0.9), RGB_BLU)
    add_rect(slide, Emu(0), Inches(0.9), SW, Inches(0.07), RGB_AZZ)
    add_text(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.65),
             title, size=24, bold=True, color=RGB_BIANCO,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.68), Inches(12.3), Inches(0.28),
                 subtitle, size=13, bold=False, color=RGB_AZZ_LT,
                 anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, n, tot):
    add_rect(slide, Emu(0), Inches(7.25), SW, Inches(0.25), RGB_GRIG_LT)
    add_text(slide, Inches(0.3), Inches(7.25), Inches(7), Inches(0.25),
             'Marco Calì — Progetto ELLONA  |  Motivazione esclusione cmos4',
             size=9, color=RGB_GRIGIO, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11), Inches(7.25), Inches(2), Inches(0.25),
             f'{n} / {tot}', size=9, color=RGB_GRIGIO,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_img(slide, path, *, top=Inches(1.1), height=Inches(5.8), max_w=Inches(12.5)):
    path = str(path)
    try:
        im = Image.open(path)
        aspect = im.size[0] / im.size[1]
        h = height
        w = Emu(int(h * aspect))
        if w > max_w:
            w = max_w
            h = Emu(int(w / aspect))
        left = Emu((SW - w) // 2)
        slide.shapes.add_picture(path, left, top, width=w, height=h)
    except Exception as e:
        print(f'  Warning: {path} → {e}')


def slide_section(num, title, sub):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SW, SH, RGB_BLU)
    add_rect(s, Emu(0), Inches(3.0), SW, Inches(0.08), RGB_AZZ_LT)
    add_text(s, Inches(1), Inches(1.8), Inches(11), Inches(0.7),
             f'Passo {num}', size=18, bold=True, color=RGB_AZZ_LT)
    add_text(s, Inches(1), Inches(2.6), Inches(11), Inches(1.4),
             title, size=38, bold=True, color=RGB_BIANCO)
    add_text(s, Inches(1), Inches(4.2), Inches(11), Inches(0.9),
             sub, size=18, color=RGB_AZZ_LT)
    return s


def slide_image_full(title, img_path, sub=None):
    s = prs.slides.add_slide(BLANK)
    add_header(s, title, sub)
    top = Inches(1.05) if not sub else Inches(1.05)
    add_img(s, img_path, top=top, height=Inches(5.85))
    return s


def slide_bullets_image(title, bullets, img_path, sub=None):
    s = prs.slides.add_slide(BLANK)
    add_header(s, title, sub)
    y0 = Inches(1.05)
    add_bullets(s, Inches(0.5), y0, Inches(5.4), Inches(5.8), bullets, size=15)
    try:
        im = Image.open(str(img_path))
        aspect = im.size[0] / im.size[1]
        h = Inches(5.0)
        w = Emu(int(h * aspect))
        max_w = Inches(6.8)
        if w > max_w:
            w = max_w
            h = Emu(int(w / aspect))
        left = SW - w - Inches(0.3)
        s.shapes.add_picture(str(img_path), left, y0 + Inches(0.2), width=w, height=h)
    except Exception as e:
        print(f'  Warning image: {e}')
    return s


def slide_bullets_only(title, bullets, sub=None):
    s = prs.slides.add_slide(BLANK)
    add_header(s, title, sub)
    add_bullets(s, Inches(0.8), Inches(1.15), Inches(11.8), Inches(5.8),
                bullets, size=17)
    return s


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITOLO
# ══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, Emu(0), Inches(4.3), SW, Inches(3.2), RGB_BLU)
add_rect(s1, Emu(0), Inches(4.2), SW, Inches(0.1), RGB_AZZ)

add_text(s1, Inches(0.8), Inches(0.5), Inches(12), Inches(0.5),
         'Politecnico di Milano — Progetto ELLONA', size=14, bold=True, color=RGB_BLU)
add_text(s1, Inches(0.8), Inches(1.5), Inches(12), Inches(1.3),
         'Esclusione di cmos4 dall\'array MOX', size=42, bold=True, color=RGB_BLU)
add_text(s1, Inches(0.8), Inches(2.9), Inches(12), Inches(0.8),
         'Motivazione empirica basata su analisi ELLONA_13 · 14 · 15',
         size=20, color=RGB_AZZ)
add_text(s1, Inches(0.8), Inches(4.55), Inches(11), Inches(0.5),
         'Progetto ELLONA  ·  Monitoraggio odori  ·  PoliMi — IREN',
         size=14, bold=True, color=RGB_BIANCO)
add_text(s1, Inches(0.8), Inches(5.3), Inches(11), Inches(0.6),
         'Marco Calì', size=24, bold=True, color=RGB_BIANCO)
add_text(s1, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
         'Aprile 2026', size=15, color=RGB_AZZ_LT)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CONTESTO: il comportamento anomalo di cmos4
# ══════════════════════════════════════════════════════════════════════════
slide_section('1', 'Il comportamento anomalo di cmos4',
              'Deriva ×59, loading negativo, CV = 100%')

slide_bullets_only(
    'Cosa sappiamo già di cmos4',
    [
        'Deriva stagionale ×59 (marzo → settembre): 20 kΩ → 1.2 MΩ  —  CV = 100%',
        'cmos1, cmos2, cmos3 hanno variazioni ×1.5 – ×5  (CV < 50%)',
        'Loading in PC₁ = −0.336: l\'unico sensore con segno negativo',
        'Effetto: in estate cmos4 alto → spinge PC₁ verso il basso → genera falsi eventi',
        'Anti-correlato con cmos2 nel baseline (r = −0.33): comportamento stagionale opposto',
        'Domanda aperta: questo è un sensore per gas ossidanti (O₃/NOₓ), non per VOC?',
    ],
    sub='Le evidenze accumulate nei mesi precedenti — riepilogo del problema'
)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — IPOTESI UMIDITÀ (ELLONA_13)
# ══════════════════════════════════════════════════════════════════════════
slide_section('2', 'Test ipotesi — Umidità assoluta',
              'ELLONA_13: power law y = A·AH^k₁ + c')

slide_bullets_image(
    'Ipotesi 1: la deriva di cmos4 è dovuta all\'umidità assoluta',
    [
        'Modello testato: y = A·AH^k₁ + c  (power law standard per MOX)',
        'Fitting con nlinfit robusto su baseline IQR weekly @1h',
        'Risultato per cmos4: R² = −0.09  (peggio di una linea piatta)',
        'AH varia da 7 a 14 g/m³ (×2) — troppo poco per spiegare ×59',
        'Timing sbagliato: AH picca ad agosto, cmos4 picca a settembre',
        'Settembre: modello predice 83 kΩ, effettivo 1320 kΩ → 93.7% non spiegato',
        '→ IPOTESI RIGETTATA: l\'umidità non spiega la deriva',
    ],
    FIG_13 / 'fig03_fit_quality.png',
    sub='R² ≈ 0 o negativo per tutti i sensori — il modello AH non è predittivo'
)

slide_image_full(
    'Scatter cmos4 vs AH — nessuna struttura power law',
    FIG_13 / 'fig04_monthly_actual_vs_predicted.png',
    sub='Effettivo vs predetto da AH per mese: settembre è il residuo più grande'
)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — IPOTESI T e RH SEPARATI (ELLONA_14)
# ══════════════════════════════════════════════════════════════════════════
slide_section('3', 'Test ipotesi — T e RH separati',
              'ELLONA_14: 5 modelli, correlazione per scala temporale')

slide_image_full(
    'Correlazione di Pearson: sensori vs variabili ambientali',
    FIG_14 / 'fig01_pearson_correlation.png',
    sub='A scala mensile: r(cmos4, T) = +0.08 ≈ 0  |  r(cmos4, RH) = +0.30  — nessuna spiega la deriva'
)

slide_bullets_only(
    'Risultato ELLONA_14: nessuna variabile disponibile spiega cmos4',
    [
        'r(cmos4, T)_mensile  = +0.079 ≈ 0  →  temperatura non spiega la deriva stagionale',
        'r(cmos4, RH)_mensile = +0.302       →  debole, non sufficiente',
        'r(cmos4, AH)_mensile = +0.244       →  ancora più debole',
        'R² di tutti i modelli (T, RH, AH, T×RH, exp(T)×RH) ≈ 0 o negativo',
        'Confronto: cmos3 ha r(T)_mensile = +0.51  →  comportamento MOX normale',
        'Il timing di cmos4 (picco settembre) non è allineato con T (picco luglio-agosto)',
        '→ cmos4 risponde a qualcosa che non misuriamo (O₃/NOₓ fotochimico?)',
    ],
    sub='La diagnosi è definitiva: cmos4 non è spiegabile con i dati ambientali disponibili'
)

slide_image_full(
    'Profili stagionali normalizzati — il timing di cmos4 è diverso',
    FIG_14 / 'fig03_seasonal_profiles.png',
    sub='cmos4 picca a settembre; T picca a luglio-agosto; RH è piatta — nessun allineamento'
)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — IMPATTO SULLA PIPELINE (ELLONA_15)
# ══════════════════════════════════════════════════════════════════════════
slide_section('4', 'Impatto sulla pipeline — 3 vs 4 sensori',
              'ELLONA_15: confronto quantitativo PCA con e senza cmos4')

slide_image_full(
    'Loading plot — 4 sensori vs 3 sensori',
    FIG_15 / 'fig01_loadings_3vs4.png',
    sub='Senza cmos4: tutti i loadings positivi, PC₁ spiega il 57.8% (+13.7pp) della varianza'
)

slide_image_full(
    'Bias stagionale di PC₁ — si riduce ma non sparisce',
    FIG_15 / 'fig02_seasonal_bias_PC1.png',
    sub='Range mensile: 2.73σ (4s) → 2.54σ (3s)  |  cmos1-3 hanno ancora variabilità stagionale residua'
)

slide_bullets_image(
    'Effetto sulla detection: −67% di eventi',
    [
        '4 sensori: 0.57% eventi sul dataset (LOD fisso k=3)',
        '3 sensori: 0.19% eventi sul dataset (−67%)',
        'Gli eventi rimossi erano falsi positivi stagionali generati da cmos4',
        'Gli eventi condivisi tra 3s e 4s: >99% concordanza → quelli reali si preservano',
        'σ_roll / σ_global: 1.11 (4s) → 1.01 (3s)',
        '→ il baseline è quasi stazionario su finestra 7 giorni senza cmos4',
    ],
    FIG_15 / 'fig06_event_concordance.png',
    sub='I falsi positivi erano generati dal loading negativo di cmos4 in estate'
)

slide_image_full(
    'Stazionarietà del baseline: σ_roll vs σ_global',
    FIG_15 / 'fig04_sigma_roll_vs_global.png',
    sub='Senza cmos4: σ_roll ≈ σ_global → la varianza stagionale esce dal baseline  |  Rolling LOD meno urgente'
)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE FINALE — CONCLUSIONE E DECISIONE
# ══════════════════════════════════════════════════════════════════════════
s_end = prs.slides.add_slide(BLANK)
add_rect(s_end, Emu(0), Emu(0), SW, SH, RGB_BLU)
add_rect(s_end, Emu(0), Inches(1.3), SW, Inches(0.06), RGB_AZZ_LT)
add_rect(s_end, Emu(0), Inches(5.9), SW, Inches(0.06), RGB_AZZ_LT)

add_text(s_end, Inches(0.8), Inches(0.3), Inches(12), Inches(0.9),
         'Conclusione e decisione', size=28, bold=True, color=RGB_BIANCO,
         anchor=MSO_ANCHOR.MIDDLE)

verdicts = [
    ('Umidità assoluta (AH)',   'R² ≈ −0.09  |  timing sbagliato',   'RIGETTATA'),
    ('Temperatura (T)',         'r_mensile = +0.08 ≈ 0',              'RIGETTATA'),
    ('Umidità relativa (RH)',   'r_mensile = +0.30  —  insufficiente','RIGETTATA'),
    ('T × RH  +  exp(T)×RH',  'R² ≈ 0 per tutti i modelli',         'RIGETTATA'),
]
y_v = Inches(1.55)
for hyp, evidence, verdict in verdicts:
    add_text(s_end, Inches(0.7), y_v, Inches(5.2), Inches(0.45),
             hyp, size=14, bold=True, color=RGB_AZZ_LT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s_end, Inches(5.9), y_v, Inches(5.0), Inches(0.45),
             evidence, size=13, color=RGB_GRIG_LT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s_end, Inches(10.9), y_v, Inches(2.1), Inches(0.45),
             verdict, size=13, bold=True, color=RGB_ROSSO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y_v += Inches(0.52)

add_text(s_end, Inches(0.8), Inches(4.05), Inches(11.7), Inches(1.7),
         'cmos4 è escluso dalla pipeline perché:\n'
         '① Il suo comportamento stagionale non è spiegabile con le variabili ambientali disponibili\n'
         '② Il suo loading negativo genera sistematicamente falsi positivi estivi\n'
         '③ La sua rimozione riduce gli eventi del 67% senza perdere quelli reali (concordanza >99%)\n'
         '④ Senza cmos4, PC₁ spiega il 57.8% della varianza (vs 44.1%) ed è più stazionario',
         size=15, color=RGB_BIANCO, italic=False)

add_text(s_end, Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.8),
         'Pipeline aggiornata: ELLONA_08_3s → ELLONA_11_3s → ELLONA_12_3s  '
         '(cartella pipeline_3sensori)',
         size=13, color=RGB_AZZ_LT, italic=True, anchor=MSO_ANCHOR.MIDDLE)

# ── FOOTER su tutte tranne prima e ultima ─────────────────────────────────
slides_list = [s for s in prs.slides]
total = len(slides_list)
for i, s in enumerate(slides_list):
    if 0 < i < total - 1:
        add_footer(s, i + 1, total)

# ── SALVA ─────────────────────────────────────────────────────────────────
out_path = OUTDIR / 'cmos4_exclusion.pptx'
prs.save(str(out_path))
print(f'\n✓ Salvato: {out_path}')
print(f'  Slide totali: {total}')
